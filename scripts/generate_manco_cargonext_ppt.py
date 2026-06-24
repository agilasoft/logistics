#!/usr/bin/env python3
"""Generate CargoNext ManCom executive presentation (PowerPoint)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Brand colours
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
TEAL = RGBColor(0x00, 0x7A, 0x87)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
MID = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
ACCENT = RGBColor(0xE8, 0x6C, 0x00)

OUT = Path(__file__).resolve().parent.parent / "docs" / "CargoNext_ManCom_Presentation.pptx"


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="CargoNext v1 Astraea  |  ManCom Briefing  |  Confidential"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MID
    p.alignment = PP_ALIGN.CENTER


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.15))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xE8)


def add_bullets(slide, items, left=0.55, top=1.45, width=8.9, height=5.3, font_size=16, level_sizes=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    level_sizes = level_sizes or {0: font_size, 1: font_size - 2}
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(level_sizes.get(level, font_size - 2))
        p.font.color.rgb = DARK if level == 0 else MID
        p.space_after = Pt(6)


def add_table_slide(slide, headers, rows, top=1.5):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, Inches(0.45), Inches(top), Inches(9.1), Inches(0.38 * nrows)).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    tbox = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "CargoNext v1 Astraea"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    for line in [
        "Integrated Logistics on ERPNext",
        "From Quote to Delivery",
        "",
        "Management Committee Briefing",
        "Agilasoft Cloud Technologies Inc.  ·  BlueCore Solutions Corp.",
    ]:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(18 if "Management" not in line else 16)
        p2.font.color.rgb = RGBColor(0xB8, 0xD4, 0xE8) if line else WHITE
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

    # --- Slide 2: Executive Summary ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Executive Summary")
    add_bullets(
        s,
        [
            "CargoNext is an integrated logistics platform on Frappe / ERPNext v16 — quote, execution, recognition, and billing in one system.",
            "v1 Astraea is Generally Available: sea, air, transport, customs, warehousing, pricing, and job costing are production-ready.",
            "Early adopters in live operations: All Transport Network, All Systems Logistics, Prime Warehouse Dynamics, Fast Prime Transport.",
            "Differentiator: finance-aware operations — Main Job vs Internal Jobs, WIP/accrual, intercompany billing, credit control.",
            "Next layer: Control Tower executive dashboards across 22 organizations (GP, pipeline, operations KPIs).",
        ],
        font_size=17,
    )
    add_footer(s)

    # --- Slide 3: Strategic Context ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Strategic Context", "Why ManCom should care")
    add_bullets(
        s,
        [
            "Fragmented tools create manual handoffs between sales, operations, and finance.",
            "Limited visibility on job profitability, milestone lead times, and credit exposure.",
            "Scaling across sea, air, transport, and warehouse requires a unified data model.",
            "CargoNext positions the group to run forwarding, trucking, and warehousing on one ERP backbone.",
            "Built in the Philippines for regional supply chains — product and delivery capability in-house.",
        ],
        font_size=17,
    )
    add_footer(s)

    # --- Slide 4: Platform Architecture ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Platform Architecture")
    add_bullets(
        s,
        [
            ("Commercial layer", 0),
            ("  Pricing Center — Sales Quotes, Change Requests, charge calculation", 1),
            ("Execution layer", 0),
            ("  Sea Freight · Air Freight · Transport · Customs · Warehousing", 1),
            ("Finance & control layer", 0),
            ("  Job Management · Revenue recognition · Intercompany · Credit management", 1),
            ("Visibility layer", 0),
            ("  Operational milestones · Document compliance · Control Tower KPIs · Customer portal", 1),
        ],
        font_size=16,
    )
    add_footer(s)

    # --- Slide 5: Module Map ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Module Map & Status")
    add_table_slide(
        s,
        ["Module", "Scope", "Status"],
        [
            ["Sea / Air Freight", "Bookings, shipments, consolidations, MBL/MAWB", "GA — Live"],
            ["Transport", "Orders, jobs, legs, run sheets, POD", "GA — Live"],
            ["Customs", "Declarations, permits, exemptions", "GA — Live"],
            ["Warehousing", "Inbound, release, VAS, contracts, gate pass", "GA — Live"],
            ["Pricing Center", "Quotes, separate billings, charge routing", "GA — Live"],
            ["Job Management", "Costing, WIP/accrual, recognition, GP views", "GA — Live"],
            ["Credit / Intercompany", "Holds, lifts, internal & IC billing", "GA — Live"],
            ["Control Tower", "22-org dashboards, pipeline, risk register", "Rollout"],
            ["MICE / High Value / Exhibits", "Specialized vertical workflows", "Selective"],
        ],
        top=1.35,
    )
    add_footer(s)

    # --- Slide 6: Quote to Cash ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Core Workflow: Quote to Cash")
    add_bullets(
        s,
        [
            "1. Sales Quote — unified charges, tax templates, Bill To / Pay To across services",
            "2. Operational documents — Sea/Air Booking, Transport Order, Declaration Order, Warehouse orders",
            "3. Execution — shipments, jobs, legs, declarations, warehouse operations",
            "4. Milestones — planned vs actual timeline per job (sea, air, transport, customs)",
            "5. Recognition — WIP and accrual per company Recognition Policy Settings",
            "6. Billing — customer invoice, internal Journal Entry, intercompany SI/PI",
            "7. Credit control — warn on save; hold create/submit/print; Credit Hold Lift Request",
        ],
        font_size=15,
    )
    add_footer(s)

    # --- Slide 7: Execution Highlights ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Execution Modules — ManCom Highlights")
    add_bullets(
        s,
        [
            ("Sea & Air", 0),
            ("Volume/weight roll-up at shipment; estimated vs actual charges for WIP and invoicing", 1),
            ("Transport", 0),
            ("Order → job → leg → run sheet → POD; dangerous goods; Lalamove / Transportify integrations", 1),
            ("Customs", 0),
            ("Permit requirements and exemptions; declaration lifecycle to clearance and release", 1),
            ("Warehousing", 0),
            ("Inbound through release, periodic billing, customer portal access", 1),
        ],
        font_size=15,
    )
    add_footer(s)

    # --- Slide 8: Finance & Governance ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Finance & Governance")
    add_bullets(
        s,
        [
            "Main Job holds customer-facing charges; Internal Jobs hold service-specific costs with revenue allocation.",
            "One Recognition Policy Settings document per company — shared date basis for WIP and accrual.",
            "Proforma GL and profitability views — expected postings from operational data.",
            "Automatic WIP/accrual reversal when real invoices and internal billing post.",
            "Credit Management: per-DocType or apply-all holds; Credit Manager approval for lifts.",
            "Intercompany: Journal Entry (same company) vs Sales/Purchase Invoice (cross-entity).",
        ],
        font_size=15,
    )
    add_footer(s)

    # --- Slide 9: Operational Milestones ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Operational Milestone Tracking")
    add_table_slide(
        s,
        ["Mode", "Key milestone chain"],
        [
            ["Sea Freight", "Booking → Gate-In → Loaded → Departed → Arrived → Delivered → Closed"],
            ["Air Freight", "Booking → Departed → In-Transit → Arrived → Delivered"],
            ["Transport", "Pick-Up → In-Transit → Delivered"],
            ["Customs", "Submitted → Under Review → Approved → Released"],
            ["Warehouse", "Start → Received → Putaway → Pick → Release → End"],
        ],
        top=1.4,
    )
    box = s.shapes.add_textbox(Inches(0.55), Inches(5.2), Inches(8.9), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "Control Tower tracks avg lead time per milestone — actionable KPI for operations meetings."
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = TEAL
    add_footer(s)

    # --- Slide 10: Control Tower ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Control Tower", "Executive visibility — 22 organizations")
    add_bullets(
        s,
        [
            ("Financial", 0),
            ("GP YTD vs target · 3-year compare · Top clients/agents/carriers · Unbilled shipments", 1),
            ("Operations", 0),
            ("Open job files · Avg age · Milestone lead times · Trips/month · Facility occupancy", 1),
            ("Pipeline & risk", 0),
            ("Exhibits, projects, luxury pipeline · Risk register entries", 1),
            ("Support functions", 0),
            ("IT tickets & uptime · HR labor vs budget · Collections 60+ aging · Credit exposure", 1),
            ("Status: seeded on install/migrate; dimension mapping configured per tenant.", 0),
        ],
        font_size=14,
    )
    add_footer(s)

    # --- Slide 11: Technology ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Technology Platform")
    add_table_slide(
        s,
        ["Component", "Requirement"],
        [
            ["Frappe Framework", "v16+"],
            ["ERPNext", "v16+"],
            ["Python", "3.14+"],
            ["Deployment", "bench install-app logistics · bench migrate"],
            ["Documentation", "docs.cargonext.com · in-app wiki user guide"],
            ["Integrations", "Lalamove, Transportify, ERPNext GL / SI / PI"],
        ],
        top=1.5,
    )
    add_footer(s)

    # --- Slide 12: Implementation Stage ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Implementation — Current Stage")
    add_table_slide(
        s,
        ["Area", "Stage", "Evidence"],
        [
            ["Product v1 Astraea", "GA — Shipped", "Release notes, press release, migration patches"],
            ["Core execution modules", "Production", "Sea, air, transport, customs, warehousing live"],
            ["Early adopter sites", "Live / expanding", "ATN, ASL, Prime Warehouse, Fast Prime"],
            ["Finance alignment", "Mature on v1", "Recognition, WIP reversal, intercompany, credit"],
            ["Control Tower", "Active rollout", "22-org registry, dashboards, pipeline/risk"],
            ["Verticals (MICE, Exhibits)", "Selective", "Enabled per business unit need"],
        ],
        top=1.35,
    )
    add_footer(s)

    # --- Slide 13: Programme Milestones ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Programme Milestones")
    add_table_slide(
        s,
        ["#", "Milestone", "Status", "Success criteria"],
        [
            ["M1", "v1 Astraea launch", "Complete", "GA release and documentation"],
            ["M2", "Early adopter go-lives", "Ongoing", "4+ named customers in production"],
            ["M3", "Finance alignment", "Complete", "Recognition, WIP reversal, IC billing live"],
            ["M4", "Quote-to-execution fidelity", "Complete", "Charges flow quote → booking/order"],
            ["M5", "Control Tower pilot", "Next", "Dashboards for priority profit centers"],
            ["M6", "Control Tower group rollout", "Planned", "All 22 orgs with live KPIs"],
            ["M7", "Portal & integration hardening", "Ongoing", "Stable customer portal and carrier APIs"],
        ],
        top=1.25,
    )
    add_footer(s)

    # --- Slide 14: Business Value ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, WHITE)
    add_title_bar(s, "Business Value for the Group")
    add_bullets(
        s,
        [
            "Revenue protection — credit holds and exposure visibility before shipment execution.",
            "Margin visibility — job-level gross profit from GL, not offline spreadsheets.",
            "Cycle time — milestone lead-time KPIs surfaced in Control Tower.",
            "Operational scale — one platform for forwarding, trucking, and warehousing.",
            "Compliance — customs permits, dangerous goods, document checklists with alerts.",
            "Strategic asset — world-class logistics software built and owned in the Philippines.",
        ],
        font_size=16,
    )
    add_footer(s)

    # --- Slide 15: Decisions / CTA ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    tbox = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.5))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "Recommended ManCom Actions"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    for item in [
        "1. Endorse CargoNext v1 Astraea as the group logistics platform standard on ERPNext.",
        "2. Approve Control Tower pilot for priority profit centers (GP, pipeline, milestone KPIs).",
        "3. Confirm rollout sequence for remaining business units (transport, warehousing, forwarding).",
        "4. Mandate post-migrate validation: recognition policy, credit rules, sample quote-to-invoice flows.",
        "",
        "Contacts: info@agilasoft.com  |  docs.cargonext.com  |  github.com/agilasoft/logistics",
    ]:
        p2 = tf.add_paragraph()
        p2.text = item
        p2.font.size = Pt(17)
        p2.font.color.rgb = RGBColor(0xD4, 0xE8, 0xF5)
        p2.space_before = Pt(10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Created: {OUT}")
    return OUT


if __name__ == "__main__":
    build()
